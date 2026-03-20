"""
OutputFormatter - Python backend implementation for generating presentations and PDFs.
Uses python-pptx for PowerPoint generation and reportlab for PDF generation.
"""

import io
from datetime import datetime
from typing import Dict, Any, Tuple, List

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    print("Warning: python-pptx not installed. Install with: pip install python-pptx")
    Presentation = None

try:
    from reportlab.lib.pagesizes import letter, A4
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import inch
    from reportlab.lib.colors import HexColor
except ImportError:
    print("Warning: reportlab not installed. Install with: pip install reportlab")
    SimpleDocTemplate = None


class OutputFormatter:
    """Handles generation of PowerPoint presentations and PDF documents from SOP content."""
    
    def __init__(self):
        self.styles = self._create_styles() if SimpleDocTemplate else None
    
    def _create_styles(self):
        """Create custom styles for PDF generation."""
        styles = getSampleStyleSheet()
        
        # Custom title style
        styles.add(ParagraphStyle(
            name='CustomTitle',
            parent=styles['Title'],
            fontSize=24,
            spaceAfter=30,
            textColor=HexColor('#1a1a2e'),
            alignment=1  # Center
        ))
        
        # Custom heading style
        styles.add(ParagraphStyle(
            name='CustomHeading',
            parent=styles['Heading1'],
            fontSize=18,
            spaceBefore=20,
            spaceAfter=12,
            textColor=HexColor('#2c3e50')
        ))
        
        # Custom step style
        styles.add(ParagraphStyle(
            name='StepStyle',
            parent=styles['Normal'],
            fontSize=12,
            spaceBefore=8,
            spaceAfter=8,
            leftIndent=20,
            borderWidth=1,
            borderColor=HexColor('#4a90d9'),
            borderPadding=10,
            backColor=HexColor('#f8f9fa')
        ))
        
        return styles
    
    def create_slide_presentation(self, content: Dict[str, Any]) -> Tuple[bytes, int]:
        """
        Generate PowerPoint presentation from content.
        
        Args:
            content: Generated content dictionary
            
        Returns:
            Tuple of (pptx_bytes, slide_count)
        """
        if not Presentation:
            raise ImportError("python-pptx is required for presentation generation")
        
        prs = Presentation()
        slide_count = 0
        
        # Title slide
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = content.get('summary', {}).get('title', 'SOP Training Presentation')
        generated_at = content.get('generatedAt', datetime.now().isoformat())
        if isinstance(generated_at, str):
            try:
                generated_at = datetime.fromisoformat(generated_at.replace('Z', '+00:00'))
            except:
                generated_at = datetime.now()
        subtitle.text = f"Generated: {generated_at.strftime('%B %d, %Y')}"
        slide_count += 1
        
        # Summary slides
        summary = content.get('summary', {})
        if summary:
            # Overview slide
            if summary.get('overview'):
                slide = self._add_content_slide(prs, "Summary", [summary['overview']])
                slide_count += 1
            
            # Key points slide
            key_points = summary.get('keyPoints', summary.get('key_points', []))
            if key_points:
                slide = self._add_content_slide(prs, "Key Points", key_points[:6])
                slide_count += 1
            
            # Safety requirements slide
            safety_reqs = summary.get('safetyRequirements', summary.get('safety_requirements', []))
            if safety_reqs:
                slide = self._add_content_slide(prs, "Safety Requirements", safety_reqs[:6])
                slide_count += 1
        
        # Training material slides
        training = content.get('trainingMaterial', content.get('training_material', {}))
        if training:
            # Learning objectives
            objectives = training.get('learningObjectives', training.get('learning_objectives', []))
            if objectives:
                slide = self._add_content_slide(prs, "Learning Objectives", objectives[:6])
                slide_count += 1
            
            # Training steps (paginated)
            steps = training.get('steps', [])
            if steps:
                for i in range(0, len(steps), 4):
                    chunk = steps[i:i+4]
                    step_texts = []
                    for j, step in enumerate(chunk):
                        if isinstance(step, dict):
                            step_text = step.get('description', step.get('title', str(step)))
                        else:
                            step_text = str(step)
                        step_texts.append(f"Step {i+j+1}: {step_text}")
                    
                    title = f"Training Steps {i+1}–{min(i+4, len(steps))}"
                    slide = self._add_content_slide(prs, title, step_texts)
                    slide_count += 1
        
        # Evaluation slide
        evaluation = content.get('evaluation', {})
        if evaluation and evaluation.get('questions'):
            questions = evaluation['questions'][:5]
            question_texts = []
            for i, q in enumerate(questions):
                if isinstance(q, dict):
                    question_text = q.get('question', str(q))
                else:
                    question_text = str(q)
                question_texts.append(f"{i+1}. {question_text}")
            
            slide = self._add_content_slide(prs, "Evaluation Questions", question_texts)
            slide_count += 1
        
        # Closing slide
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "Thank You"
        subtitle.text = "Please complete the evaluation to confirm your understanding."
        slide_count += 1
        
        # Save to bytes
        pptx_buffer = io.BytesIO()
        prs.save(pptx_buffer)
        pptx_buffer.seek(0)
        
        return pptx_buffer.getvalue(), slide_count
    
    def _add_content_slide(self, prs: 'Presentation', title: str, bullets: List[str]) -> Any:
        """Add a content slide with title and bullet points."""
        slide_layout = prs.slide_layouts[1]  # Title and Content layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Set title
        title_shape = slide.shapes.title
        title_shape.text = title
        
        # Add bullet points
        content_placeholder = slide.placeholders[1]
        text_frame = content_placeholder.text_frame
        text_frame.clear()
        
        for i, bullet in enumerate(bullets):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = bullet
            p.level = 0
            p.font.size = Pt(16)
        
        return slide
    
    def generate_pdf(self, content: Dict[str, Any]) -> Tuple[bytes, int]:
        """
        Generate PDF document from content.
        
        Args:
            content: Generated content dictionary
            
        Returns:
            Tuple of (pdf_bytes, page_count)
        """
        if not SimpleDocTemplate:
            raise ImportError("reportlab is required for PDF generation")
        
        buffer = io.BytesIO()
        doc = SimpleDocTemplate(buffer, pagesize=letter, topMargin=1*inch)
        story = []
        
        # Title
        title = content.get('summary', {}).get('title', 'SOP Training Material')
        story.append(Paragraph(title, self.styles['CustomTitle']))
        
        # Generated date
        generated_at = content.get('generatedAt', datetime.now().isoformat())
        if isinstance(generated_at, str):
            try:
                generated_at = datetime.fromisoformat(generated_at.replace('Z', '+00:00'))
            except:
                generated_at = datetime.now()
        
        story.append(Paragraph(f"Generated: {generated_at.strftime('%B %d, %Y at %I:%M %p')}", 
                              self.styles['Normal']))
        story.append(Spacer(1, 20))
        
        # Summary section
        summary = content.get('summary', {})
        if summary:
            story.append(Paragraph("Summary", self.styles['CustomHeading']))
            
            if summary.get('overview'):
                story.append(Paragraph(summary['overview'], self.styles['Normal']))
                story.append(Spacer(1, 12))
            
            # Key points
            key_points = summary.get('keyPoints', summary.get('key_points', []))
            if key_points:
                story.append(Paragraph("Key Points", self.styles['Heading2']))
                for point in key_points:
                    story.append(Paragraph(f"• {point}", self.styles['Normal']))
                story.append(Spacer(1, 12))
            
            # Safety requirements
            safety_reqs = summary.get('safetyRequirements', summary.get('safety_requirements', []))
            if safety_reqs:
                story.append(Paragraph("Safety Requirements", self.styles['Heading2']))
                for req in safety_reqs:
                    story.append(Paragraph(f"• {req}", self.styles['Normal']))
                story.append(Spacer(1, 12))
        
        # Training material section
        training = content.get('trainingMaterial', content.get('training_material', {}))
        if training:
            story.append(PageBreak())
            story.append(Paragraph("Training Material", self.styles['CustomHeading']))
            
            # Learning objectives
            objectives = training.get('learningObjectives', training.get('learning_objectives', []))
            if objectives:
                story.append(Paragraph("Learning Objectives", self.styles['Heading2']))
                for obj in objectives:
                    story.append(Paragraph(f"• {obj}", self.styles['Normal']))
                story.append(Spacer(1, 12))
            
            # Training steps
            steps = training.get('steps', [])
            if steps:
                story.append(Paragraph("Training Steps", self.styles['Heading2']))
                for i, step in enumerate(steps):
                    if isinstance(step, dict):
                        step_text = step.get('description', step.get('title', str(step)))
                        step_num = step.get('step_number', i + 1)
                    else:
                        step_text = str(step)
                        step_num = i + 1
                    
                    story.append(Paragraph(f"<b>Step {step_num}:</b> {step_text}", 
                                         self.styles['StepStyle']))
                story.append(Spacer(1, 12))
            
            # Duration
            duration = training.get('estimatedDuration', training.get('estimated_duration'))
            if duration:
                story.append(Paragraph(f"<i>Estimated duration: {duration} minutes</i>", 
                                     self.styles['Normal']))
        
        # Evaluation section
        evaluation = content.get('evaluation', {})
        if evaluation:
            story.append(PageBreak())
            story.append(Paragraph("Evaluation", self.styles['CustomHeading']))
            
            if evaluation.get('instructions'):
                story.append(Paragraph(evaluation['instructions'], self.styles['Normal']))
                story.append(Spacer(1, 12))
            
            questions = evaluation.get('questions', [])
            if questions:
                story.append(Paragraph("Questions", self.styles['Heading2']))
                for i, q in enumerate(questions):
                    if isinstance(q, dict):
                        question_text = q.get('question', str(q))
                    else:
                        question_text = str(q)
                    
                    story.append(Paragraph(f"<b>Q{i+1}:</b> {question_text}", 
                                         self.styles['Normal']))
                    story.append(Spacer(1, 8))
            
            passing_score = evaluation.get('passingScore', evaluation.get('passing_score'))
            if passing_score:
                story.append(Spacer(1, 12))
                story.append(Paragraph(f"<i>Passing score: {passing_score}%</i>", 
                                     self.styles['Normal']))
        
        # Build PDF
        doc.build(story)
        buffer.seek(0)
        
        # Estimate page count (rough calculation)
        page_count = max(1, len(story) // 20)  # Rough estimate
        
        return buffer.getvalue(), page_count